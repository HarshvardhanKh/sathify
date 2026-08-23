import asyncio

from pptx import Presentation
from pypdf import PdfWriter

from app.models.document import EquationBlockModel, HeadingBlock, ParagraphBlock
from app.services.documents.pdf_parser import extract_pdf_pages
from app.services.documents.semantic_converter import _json_to_blocks, _naive_paragraph_blocks, _strip_fences, convert_raw_text
from app.services.documents.slides_parser import extract_slides
from app.services.documents.store import DocumentStore
from app.models.document import Document


def run(coro):
    return asyncio.run(coro)


def test_extract_pdf_pages_real_file(tmp_path):
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    pdf_path = tmp_path / "test.pdf"
    with open(pdf_path, "wb") as f:
        writer.write(f)

    pages = extract_pdf_pages(pdf_path)
    assert len(pages) == 2  # correct page count even with no text content


def test_extract_slides_real_pptx(tmp_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # title + content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Neural Networks"
    body = slide.placeholders[1]
    body.text_frame.text = "Input layer"
    body.text_frame.add_paragraph().text = "Hidden layers"

    pptx_path = tmp_path / "test.pptx"
    prs.save(pptx_path)

    slides = extract_slides(pptx_path)
    assert len(slides) == 1
    assert slides[0].title == "Neural Networks"
    assert "Input layer" in slides[0].bullets
    assert "Hidden layers" in slides[0].bullets


def test_naive_paragraph_blocks_splits_on_blank_lines():
    text = "First paragraph.\n\nSecond paragraph.\n\nThird."
    blocks = _naive_paragraph_blocks(text)
    assert len(blocks) == 3
    assert all(isinstance(b, ParagraphBlock) for b in blocks)
    assert blocks[0].text == "First paragraph."


def test_json_to_blocks_produces_headings_paragraphs_equations():
    parsed = {
        "sections": [{"heading": "Introduction", "content": "Some text here."}],
        "equations": [{"latex": "E = mc^2", "spoken": "E equals m c squared"}],
        "images": [{"description": "A graph showing exponential growth"}],
    }
    blocks = _json_to_blocks(parsed)
    kinds = [b.kind for b in blocks]
    assert "heading" in kinds
    assert "paragraph" in kinds
    assert "equation" in kinds
    assert "image" in kinds
    eq = next(b for b in blocks if isinstance(b, EquationBlockModel))
    assert eq.latex == "E = mc^2"


def test_strip_fences_removes_markdown_code_block():
    wrapped = '```json\n{"title": "x"}\n```'
    assert _strip_fences(wrapped) == '{"title": "x"}'
    assert _strip_fences('{"title": "x"}') == '{"title": "x"}'


def test_convert_raw_text_falls_back_honestly_without_gemini_key():
    """No GEMINI_API_KEYS configured in this test environment -> must fall
    back to naive splitting and say so, never fabricate structure."""
    blocks, title, summary, ai_used = run(convert_raw_text("Paragraph one.\n\nParagraph two."))
    assert ai_used is False
    assert "unavailable" in summary.lower()
    assert len(blocks) == 2


def test_document_store_crud_and_camel_case_serialization():
    store = DocumentStore()
    doc = Document(title="My Doc", original_format="pdf")
    store.add(doc)

    fetched = store.get(doc.id)
    assert fetched is not None
    assert fetched.title == "My Doc"
    assert doc in store.list()

    dumped = doc.model_dump(by_alias=True)
    assert "originalFormat" in dumped  # camelCase alias matches new-frontend's TS types
    assert "original_format" not in dumped
