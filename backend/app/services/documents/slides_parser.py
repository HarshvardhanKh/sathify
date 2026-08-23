"""
PPTX slide extraction via python-pptx. Produces one structured record per
slide: title, bullet outline, and speaker notes — real extraction, not
just raw text-box dumping (Phase 14's "not simply extracting text boxes"
requirement is addressed at the semantic_converter stage, which turns this
into the richer {title, outline, visual_description} shape via Gemini when
available).
"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation

logger = logging.getLogger("isl.documents.slides_parser")


@dataclass
class SlideContent:
    slide_number: int
    title: str
    bullets: list[str] = field(default_factory=list)
    notes: str = ""
    has_image: bool = False


def extract_slides(path: str | Path) -> list[SlideContent]:
    prs = Presentation(str(path))
    slides: list[SlideContent] = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets: list[str] = []
        has_image = False

        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                has_image = True
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                bullets.extend(line.strip() for line in text.split("\n") if line.strip())

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(SlideContent(
            slide_number=i, title=title or f"Slide {i}", bullets=bullets, notes=notes, has_image=has_image,
        ))

    return slides
